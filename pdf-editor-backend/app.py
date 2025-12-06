"""
PDF Editor Service - Native PDF Editing with PyMuPDF + PaddleOCR
Google Cloud Run service for real-time native PDF editing
Features: Text add/edit/delete, OCR, page rendering, export to Word/Excel/PPT
"""

from fastapi import FastAPI, UploadFile, File, HTTPException, Form, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, StreamingResponse
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import fitz  # PyMuPDF
import io
import base64
import os
import logging
import time
import uuid
import json
from datetime import datetime
import requests

# PaddleOCR imports
try:
    from paddleocr import PaddleOCR
    PADDLEOCR_AVAILABLE = True
except ImportError:
    PADDLEOCR_AVAILABLE = False
    logging.warning("PaddleOCR not available, OCR features will be disabled")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="PDF Editor Service", version="1.0.0")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory session storage (for production, use Redis or Cloud Storage)
sessions: Dict[str, Dict[str, Any]] = {}

# Initialize PaddleOCR (lazy loading)
ocr_engine = None

def get_ocr_engine():
    """Get or initialize PaddleOCR engine"""
    global ocr_engine
    if ocr_engine is None and PADDLEOCR_AVAILABLE:
        logger.info("Initializing PaddleOCR engine...")
        try:
            # Initialize with Hindi and English support
            ocr_engine = PaddleOCR(use_angle_cls=True, lang='en', use_gpu=False)
            logger.info("PaddleOCR initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize PaddleOCR: {e}")
            ocr_engine = None
    return ocr_engine

# Credit system integration
def check_user_credits(user_id: str, required_credits: int) -> Dict[str, Any]:
    """Check if user has sufficient credits"""
    try:
        # Call credit API
        api_base = os.environ.get('API_BASE_URL', 'https://easyjpgtopdf.com')
        response = requests.get(
            f"{api_base}/api/credits/balance",
            params={"userId": user_id},
            timeout=5
        )
        if response.status_code == 200:
            data = response.json()
            credits = data.get('credits', 0)
            unlimited = data.get('unlimited', False)
            is_premium = data.get('isPremium', False)
            
            if unlimited or is_premium:
                return {"has_credits": True, "unlimited": True, "credits": credits}
            
            if credits >= required_credits:
                return {"has_credits": True, "unlimited": False, "credits": credits}
            else:
                return {"has_credits": False, "unlimited": False, "credits": credits}
    except Exception as e:
        logger.error(f"Error checking credits: {e}")
    
    # Default: allow if API unavailable (for development)
    return {"has_credits": True, "unlimited": False, "credits": 0}

def deduct_credits(user_id: str, amount: int, reason: str, metadata: Dict = None):
    """Deduct credits from user account"""
    try:
        api_base = os.environ.get('API_BASE_URL', 'https://easyjpgtopdf.com')
        response = requests.post(
            f"{api_base}/api/credits/deduct",
            json={
                "userId": user_id,
                "amount": amount,
                "reason": reason,
                "metadata": metadata or {}
            },
            timeout=5
        )
        return response.status_code == 200
    except Exception as e:
        logger.error(f"Error deducting credits: {e}")
        return False

# Device fingerprint tracking for daily page limits
device_daily_usage: Dict[str, Dict[str, Any]] = {}  # {device_id: {date: page_count}}

def check_daily_page_limit(device_id: str, user_id: Optional[str] = None) -> Dict[str, Any]:
    """Check daily page edit limit (5 pages per day per device)"""
    FREE_PAGES_PER_DAY = 5
    today = datetime.now().date().isoformat()
    
    # Check if user is premium
    if user_id:
        credit_info = check_user_credits(user_id, 0)
        if credit_info.get("unlimited", False) or credit_info.get("isPremium", False):
            return {"allowed": True, "remaining": -1, "limit": -1, "is_premium": True}
    
    # Track by device fingerprint
    if device_id not in device_daily_usage:
        device_daily_usage[device_id] = {}
    
    device_usage = device_daily_usage[device_id]
    
    # Reset if new day
    if today not in device_usage:
        device_usage[today] = 0
    
    pages_used = device_usage[today]
    remaining = FREE_PAGES_PER_DAY - pages_used
    
    if pages_used < FREE_PAGES_PER_DAY:
        return {
            "allowed": True,
            "remaining": remaining,
            "limit": FREE_PAGES_PER_DAY,
            "pages_used": pages_used,
            "is_premium": False
        }
    else:
        return {
            "allowed": False,
            "remaining": 0,
            "limit": FREE_PAGES_PER_DAY,
            "pages_used": pages_used,
            "is_premium": False,
            "message": f"Daily limit reached. You've edited {FREE_PAGES_PER_DAY} pages today. Upgrade to premium for unlimited editing."
        }

def increment_daily_page_count(device_id: str):
    """Increment daily page edit count for device"""
    today = datetime.now().date().isoformat()
    if device_id not in device_daily_usage:
        device_daily_usage[device_id] = {}
    if today not in device_daily_usage[device_id]:
        device_daily_usage[device_id][today] = 0
    device_daily_usage[device_id][today] += 1

# Request models
class SessionStartRequest(BaseModel):
    userId: Optional[str] = None

class PageRenderRequest(BaseModel):
    session_id: str
    page_number: int
    zoom: float = 1.5

class TextSearchRequest(BaseModel):
    session_id: str
    query: str

class TextEditRequest(BaseModel):
    session_id: str
    page_number: int
    bbox: Optional[List[float]] = None
    match_id: Optional[str] = None
    new_text: str
    font_name: str = "Helvetica"
    font_size: float = 12
    color: List[float] = [0, 0, 0]
    userId: Optional[str] = None
    deviceId: Optional[str] = None

class TextAddRequest(BaseModel):
    session_id: str
    page_number: int
    x: float
    y: float
    text: str
    font_name: str = "Helvetica"
    font_size: float = 12
    color: List[float] = [0, 0, 0]
    userId: Optional[str] = None
    deviceId: Optional[str] = None

class TextDeleteRequest(BaseModel):
    session_id: str
    page_number: int
    bbox: Optional[List[float]] = None
    match_id: Optional[str] = None
    userId: Optional[str] = None
    deviceId: Optional[str] = None

class OCRRequest(BaseModel):
    session_id: str
    page_number: int
    userId: Optional[str] = None

class ExportRequest(BaseModel):
    session_id: str
    format: str  # pdf, docx, xlsx, pptx
    userId: Optional[str] = None

@app.get("/health")
async def health():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "service": "pdf-editor",
        "paddleocr": PADDLEOCR_AVAILABLE,
        "ocr_initialized": ocr_engine is not None if PADDLEOCR_AVAILABLE else False
    }

@app.post("/session/start")
async def start_session(file: UploadFile = File(...), userId: Optional[str] = Form(None), deviceId: Optional[str] = Form(None)):
    """Start a new PDF editing session"""
    try:
        # Read PDF file
        pdf_bytes = await file.read()
        
        # Open PDF with PyMuPDF
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        page_count = len(pdf_doc)
        
        # Generate session ID
        session_id = str(uuid.uuid4())
        
        # Store session in memory
        sessions[session_id] = {
            "pdf_doc": pdf_doc,
            "pdf_bytes": pdf_bytes,
            "page_count": page_count,
            "created_at": datetime.now().isoformat(),
            "userId": userId,
            "deviceId": deviceId,
            "edits": [],
            "pages_edited": set()  # Track which pages have been edited
        }
        
        logger.info(f"Session started: {session_id}, pages: {page_count}, user: {userId}, device: {deviceId}")
        
        return {
            "success": True,
            "session_id": session_id,
            "page_count": page_count
        }
    except Exception as e:
        logger.error(f"Error starting session: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/page/render")
async def render_page(request: PageRenderRequest):
    """Render a PDF page as PNG image"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        page_num = request.page_number - 1  # Convert to 0-based
        
        if page_num < 0 or page_num >= len(pdf_doc):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        page = pdf_doc[page_num]
        
        # Render page to image
        zoom = request.zoom
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        # Convert to PNG bytes
        img_bytes = pix.tobytes("png")
        
        return Response(content=img_bytes, media_type="image/png")
    except Exception as e:
        logger.error(f"Error rendering page: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/text/search")
async def search_text(request: TextSearchRequest):
    """Search for text in PDF"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        matches = []
        
        # Search across all pages
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            text_instances = page.search_for(request.query)
            
            for inst in text_instances:
                matches.append({
                    "page_number": page_num + 1,
                    "bbox": list(inst),
                    "text": request.query,
                    "match_id": f"{page_num}_{len(matches)}"
                })
        
        return {
            "success": True,
            "matches": matches,
            "count": len(matches)
        }
    except Exception as e:
        logger.error(f"Error searching text: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/text/edit")
async def edit_text(request: TextEditRequest):
    """Edit text in PDF (replace existing text)"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        page_num = request.page_number - 1
        
        if page_num < 0 or page_num >= len(pdf_doc):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        page = pdf_doc[page_num]
        
        # Check daily page limit by device fingerprint
        device_id = request.deviceId or "anonymous"
        page_already_edited = request.page_number in session.get("pages_edited", set())
        
        if not page_already_edited:
            limit_check = check_daily_page_limit(device_id, request.userId)
            
            if not limit_check["allowed"]:
                raise HTTPException(
                    status_code=403,
                    detail=limit_check.get("message", "Daily limit reached. Upgrade to premium for unlimited editing.")
                )
            
            # Check if we need to charge credits (after free limit)
            if limit_check["remaining"] <= 0:
                # Free limit exhausted, charge credits
                if request.userId:
                    credit_info = check_user_credits(request.userId, 10)
                    if not credit_info.get("unlimited", False):
                        if credit_info["credits"] < 10:
                            raise HTTPException(
                                status_code=403,
                                detail="Insufficient credits. 10 credits required per page edit after free limit. Please recharge your account."
                            )
                        # Deduct credits
                        deduct_credits(request.userId, 10, "PDF page edit (after free limit)", {
                            "session_id": request.session_id,
                            "page": request.page_number,
                            "action": "edit_text",
                            "device_id": device_id
                        })
                else:
                    raise HTTPException(
                        status_code=403,
                        detail="Free daily limit reached. Please login and purchase credits to continue editing."
                    )
            
            # Mark page as edited and increment daily count
            if "pages_edited" not in session:
                session["pages_edited"] = set()
            session["pages_edited"].add(request.page_number)
            increment_daily_page_count(device_id)
        
        # Find text to replace
        if request.bbox:
            bbox = fitz.Rect(request.bbox)
            # Delete old text by drawing white rectangle
            page.draw_rect(bbox, color=(1, 1, 1), fill=(1, 1, 1))
            # Add new text at same position
            point = fitz.Point(bbox.x0, bbox.y1)
        elif request.match_id:
            # Find match by ID (simplified - in production, store matches)
            text_instances = page.search_for(request.new_text)
            if text_instances:
                bbox = text_instances[0]
                page.draw_rect(bbox, color=(1, 1, 1), fill=(1, 1, 1))
                point = fitz.Point(bbox.x0, bbox.y1)
            else:
                raise HTTPException(status_code=404, detail="Text match not found")
        else:
            raise HTTPException(status_code=400, detail="bbox or match_id required")
        
        # Insert new text
        font = fitz.Font(request.font_name)
        color = tuple(c / 255.0 for c in request.color[:3])
        page.insert_text(
            point,
            request.new_text,
            fontsize=request.font_size,
            fontname=request.font_name,
            color=color
        )
        
        # Update PDF in memory (save to bytes buffer)
        pdf_bytes = pdf_doc.tobytes()
        session["pdf_bytes"] = pdf_bytes
        # Reopen to keep session active
        session["pdf_doc"] = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Track edit
        session["edits"].append({
            "type": "edit",
            "page": request.page_number,
            "timestamp": datetime.now().isoformat(),
            "device_id": device_id
        })
        
        # Get current limit status
        limit_status = check_daily_page_limit(device_id, request.userId)
        credits_used = 0
        if limit_status["remaining"] <= 0 and request.userId:
            credits_used = 10
        
        return {
            "success": True,
            "message": "Text edited successfully",
            "credits_used": credits_used,
            "daily_limit": {
                "remaining": limit_status["remaining"],
                "limit": limit_status["limit"],
                "pages_used": limit_status["pages_used"]
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error editing text: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/text/add")
async def add_text(request: TextAddRequest):
    """Add new text to PDF"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        page_num = request.page_number - 1
        
        if page_num < 0 or page_num >= len(pdf_doc):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        page = pdf_doc[page_num]
        
        # Check daily page limit by device fingerprint
        device_id = request.deviceId or "anonymous"
        page_already_edited = request.page_number in session.get("pages_edited", set())
        
        if not page_already_edited:
            limit_check = check_daily_page_limit(device_id, request.userId)
            
            if not limit_check["allowed"]:
                raise HTTPException(
                    status_code=403,
                    detail=limit_check.get("message", "Daily limit reached. Upgrade to premium for unlimited editing.")
                )
            
            # Check if we need to charge credits (after free limit)
            if limit_check["remaining"] <= 0:
                # Free limit exhausted, charge credits
                if request.userId:
                    credit_info = check_user_credits(request.userId, 10)
                    if not credit_info.get("unlimited", False):
                        if credit_info["credits"] < 10:
                            raise HTTPException(
                                status_code=403,
                                detail="Insufficient credits. 10 credits required per page edit after free limit. Please recharge your account."
                            )
                        # Deduct credits
                        deduct_credits(request.userId, 10, "PDF page edit (after free limit)", {
                            "session_id": request.session_id,
                            "page": request.page_number,
                            "action": "add_text",
                            "device_id": device_id
                        })
                else:
                    raise HTTPException(
                        status_code=403,
                        detail="Free daily limit reached. Please login and purchase credits to continue editing."
                    )
            
            # Mark page as edited and increment daily count
            if "pages_edited" not in session:
                session["pages_edited"] = set()
            session["pages_edited"].add(request.page_number)
            increment_daily_page_count(device_id)
        
        # Convert coordinates (PDF uses bottom-left origin)
        page_height = page.rect.height
        pdf_y = page_height - request.y
        
        # Insert text
        font = fitz.Font(request.font_name)
        color = tuple(c / 255.0 for c in request.color[:3])
        point = fitz.Point(request.x, pdf_y)
        
        page.insert_text(
            point,
            request.text,
            fontsize=request.font_size,
            fontname=request.font_name,
            color=color
        )
        
        # Update PDF in memory (save to bytes buffer)
        pdf_bytes = pdf_doc.tobytes()
        session["pdf_bytes"] = pdf_bytes
        # Reopen to keep session active
        session["pdf_doc"] = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Track edit
        session["edits"].append({
            "type": "add",
            "page": request.page_number,
            "timestamp": datetime.now().isoformat(),
            "device_id": device_id
        })
        
        # Get current limit status
        limit_status = check_daily_page_limit(device_id, request.userId)
        credits_used = 0
        if limit_status["remaining"] <= 0 and request.userId:
            credits_used = 10
        
        return {
            "success": True,
            "message": "Text added successfully",
            "credits_used": credits_used,
            "daily_limit": {
                "remaining": limit_status["remaining"],
                "limit": limit_status["limit"],
                "pages_used": limit_status["pages_used"]
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error adding text: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/text/delete")
async def delete_text(request: TextDeleteRequest):
    """Delete text from PDF"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        page_num = request.page_number - 1
        
        if page_num < 0 or page_num >= len(pdf_doc):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        page = pdf_doc[page_num]
        
        # Check daily page limit by device fingerprint
        device_id = request.deviceId or "anonymous"
        page_already_edited = request.page_number in session.get("pages_edited", set())
        
        if not page_already_edited:
            limit_check = check_daily_page_limit(device_id, request.userId)
            
            if not limit_check["allowed"]:
                raise HTTPException(
                    status_code=403,
                    detail=limit_check.get("message", "Daily limit reached. Upgrade to premium for unlimited editing.")
                )
            
            # Check if we need to charge credits (after free limit)
            if limit_check["remaining"] <= 0:
                # Free limit exhausted, charge credits
                if request.userId:
                    credit_info = check_user_credits(request.userId, 10)
                    if not credit_info.get("unlimited", False):
                        if credit_info["credits"] < 10:
                            raise HTTPException(
                                status_code=403,
                                detail="Insufficient credits. 10 credits required per page edit after free limit. Please recharge your account."
                            )
                        # Deduct credits
                        deduct_credits(request.userId, 10, "PDF page edit (after free limit)", {
                            "session_id": request.session_id,
                            "page": request.page_number,
                            "action": "delete_text",
                            "device_id": device_id
                        })
                else:
                    raise HTTPException(
                        status_code=403,
                        detail="Free daily limit reached. Please login and purchase credits to continue editing."
                    )
            
            # Mark page as edited and increment daily count
            if "pages_edited" not in session:
                session["pages_edited"] = set()
            session["pages_edited"].add(request.page_number)
            increment_daily_page_count(device_id)
        
        # Delete text by drawing white rectangle
        if request.bbox:
            bbox = fitz.Rect(request.bbox)
            page.draw_rect(bbox, color=(1, 1, 1), fill=(1, 1, 1))
        elif request.match_id:
            # Find and delete by match_id (simplified)
            # In production, maintain a match registry
            raise HTTPException(status_code=400, detail="match_id deletion requires bbox")
        else:
            raise HTTPException(status_code=400, detail="bbox or match_id required")
        
        # Update PDF in memory (save to bytes buffer)
        pdf_bytes = pdf_doc.tobytes()
        session["pdf_bytes"] = pdf_bytes
        # Reopen to keep session active
        session["pdf_doc"] = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Track edit
        session["edits"].append({
            "type": "delete",
            "page": request.page_number,
            "timestamp": datetime.now().isoformat(),
            "device_id": device_id
        })
        
        # Get current limit status
        limit_status = check_daily_page_limit(device_id, request.userId)
        credits_used = 0
        if limit_status["remaining"] <= 0 and request.userId:
            credits_used = 10
        
        return {
            "success": True,
            "message": "Text deleted successfully",
            "credits_used": credits_used,
            "daily_limit": {
                "remaining": limit_status["remaining"],
                "limit": limit_status["limit"],
                "pages_used": limit_status["pages_used"]
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting text: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ocr/page")
async def ocr_page(request: OCRRequest):
    """Run OCR on a PDF page"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        page_num = request.page_number - 1
        
        if page_num < 0 or page_num >= len(pdf_doc):
            raise HTTPException(status_code=400, detail="Invalid page number")
        
        if not PADDLEOCR_AVAILABLE:
            raise HTTPException(status_code=503, detail="OCR not available")
        
        ocr = get_ocr_engine()
        if not ocr:
            raise HTTPException(status_code=503, detail="OCR engine not initialized")
        
        page = pdf_doc[page_num]
        
        # Render page to image
        mat = fitz.Matrix(2, 2)  # 2x zoom for better OCR
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")
        
        # Run OCR
        import numpy as np
        from PIL import Image
        img = Image.open(io.BytesIO(img_bytes))
        img_array = np.array(img)
        
        result = ocr.ocr(img_array, cls=True)
        
        # Parse OCR results
        ocr_texts = []
        for line in result[0] if result else []:
            bbox = line[0]
            text_info = line[1]
            text = text_info[0]
            confidence = text_info[1]
            
            # Convert bbox to PDF coordinates
            x0 = min(p[0] for p in bbox) / 2  # Divide by zoom
            y0 = min(p[1] for p in bbox) / 2
            x1 = max(p[0] for p in bbox) / 2
            y1 = max(p[1] for p in bbox) / 2
            
            # Convert to PDF coordinate system (bottom-left origin)
            page_height = page.rect.height
            pdf_y0 = page_height - y1
            pdf_y1 = page_height - y0
            
            ocr_texts.append({
                "text": text,
                "bbox": [x0, pdf_y0, x1, pdf_y1],
                "confidence": confidence
            })
        
        return {
            "success": True,
            "texts": ocr_texts,
            "count": len(ocr_texts)
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error running OCR: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/export")
async def export_pdf(request: ExportRequest):
    """Export PDF to different formats"""
    try:
        session = sessions.get(request.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        
        pdf_doc = session["pdf_doc"]
        
        if request.format == "pdf":
            # Return PDF as-is
            pdf_bytes = pdf_doc.tobytes()
            return Response(
                content=pdf_bytes,
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename=edited.pdf"}
            )
        elif request.format == "docx":
            # Convert to DOCX using python-docx
            try:
                from docx import Document
                from docx.shared import Inches, Pt
                
                doc = Document()
                
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc[page_num]
                    text = page.get_text()
                    
                    # Add page break (except first page)
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # Add text
                    para = doc.add_paragraph(text)
                    para.style.font.size = Pt(12)
                
                # Save to bytes
                docx_bytes = io.BytesIO()
                doc.save(docx_bytes)
                docx_bytes.seek(0)
                
                return Response(
                    content=docx_bytes.read(),
                    media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    headers={"Content-Disposition": "attachment; filename=exported.docx"}
                )
            except ImportError:
                raise HTTPException(status_code=503, detail="DOCX export not available (python-docx not installed)")
        elif request.format == "xlsx":
            # Convert to XLSX using openpyxl
            try:
                from openpyxl import Workbook
                
                wb = Workbook()
                ws = wb.active
                ws.title = "PDF Content"
                
                row = 1
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc[page_num]
                    text = page.get_text()
                    lines = text.split('\n')
                    
                    for line in lines:
                        if line.strip():
                            ws.cell(row=row, column=1, value=line)
                            row += 1
                
                # Save to bytes
                xlsx_bytes = io.BytesIO()
                wb.save(xlsx_bytes)
                xlsx_bytes.seek(0)
                
                return Response(
                    content=xlsx_bytes.read(),
                    media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    headers={"Content-Disposition": "attachment; filename=exported.xlsx"}
                )
            except ImportError:
                raise HTTPException(status_code=503, detail="XLSX export not available (openpyxl not installed)")
        elif request.format == "pptx":
            # Convert to PPTX using python-pptx
            try:
                from pptx import Presentation
                from pptx.util import Inches
                
                prs = Presentation()
                
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc[page_num]
                    text = page.get_text()
                    
                    # Add slide
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                    
                    # Add text box
                    left = Inches(1)
                    top = Inches(1)
                    width = Inches(8)
                    height = Inches(5.5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = text
                
                # Save to bytes
                pptx_bytes = io.BytesIO()
                prs.save(pptx_bytes)
                pptx_bytes.seek(0)
                
                return Response(
                    content=pptx_bytes.read(),
                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition": "attachment; filename=exported.pptx"}
                )
            except ImportError:
                raise HTTPException(status_code=503, detail="PPTX export not available (python-pptx not installed)")
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error exporting PDF: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8080))
    uvicorn.run(app, host="0.0.0.0", port=port)

